"""
Generate all downloadable workfiles for the AI Apps specialization toolkit.

Outputs land under public/templates/{engagement,deliverables,audit}/.

Run from repo root:
    python scripts/workfiles/generate_all.py
"""
from __future__ import annotations

import os
from pathlib import Path

from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.worksheet.datavalidation import DataValidation
from openpyxl.utils import get_column_letter

REPO_ROOT = Path(__file__).resolve().parents[2]
OUT = REPO_ROOT / "public" / "templates"

BRAND_BLUE = RGBColor(0x00, 0x4B, 0x87)
BRAND_BLUE_PPT = PRGBColor(0x00, 0x4B, 0x87)
HEADER_FILL = PatternFill("solid", fgColor="004B87")
HEADER_FONT = Font(bold=True, color="FFFFFF", size=11)
STATUS_OPTIONS = '"Not started,In progress,Blocked,Complete"'
THIN = Side(border_style="thin", color="BFBFBF")
BORDER = Border(left=THIN, right=THIN, top=THIN, bottom=THIN)


# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────

def ensure_dirs() -> None:
    for sub in ("engagement", "deliverables", "audit"):
        (OUT / sub).mkdir(parents=True, exist_ok=True)


def docx_set_base_style(doc: Document) -> None:
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)


def docx_cover(doc: Document, title: str, subtitle: str, owner: str = "Practice Lead") -> None:
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("\n\n\n" + title)
    run.bold = True
    run.font.size = Pt(28)
    run.font.color.rgb = BRAND_BLUE

    p2 = doc.add_paragraph()
    p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = p2.add_run(subtitle)
    r2.italic = True
    r2.font.size = Pt(14)

    doc.add_paragraph().add_run("\n\n\n\n").italic = True
    table = doc.add_table(rows=4, cols=2)
    table.style = "Light Grid Accent 1"
    rows = [
        ("Document owner", owner),
        ("Version", "1.0 (template)"),
        ("Last updated", "yyyy-mm-dd"),
        ("Status", "Draft"),
    ]
    for (k, v), row in zip(rows, table.rows):
        row.cells[0].text = k
        row.cells[1].text = v
    doc.add_page_break()


def docx_toc(doc: Document, sections: list[str]) -> None:
    h = doc.add_heading("Table of contents", level=1)
    for run in h.runs:
        run.font.color.rgb = BRAND_BLUE
    for i, s in enumerate(sections, 1):
        doc.add_paragraph(f"{i}. {s}")
    doc.add_page_break()


def docx_section(doc: Document, title: str, prompts: list[str], extra: list[str] | None = None) -> None:
    h = doc.add_heading(title, level=1)
    for run in h.runs:
        run.font.color.rgb = BRAND_BLUE
    for prompt in prompts:
        p = doc.add_paragraph()
        r = p.add_run(prompt)
        r.italic = True
        r.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        doc.add_paragraph("")  # space for writer
    if extra:
        for line in extra:
            doc.add_paragraph(line)
    doc.add_paragraph("")


def write_word(
    path: Path,
    title: str,
    subtitle: str,
    sections: list[tuple[str, list[str], list[str] | None]],
    owner: str = "Practice Lead",
) -> None:
    doc = Document()
    docx_set_base_style(doc)
    docx_cover(doc, title, subtitle, owner)
    docx_toc(doc, [s[0] for s in sections])
    for title_s, prompts, extra in sections:
        docx_section(doc, title_s, prompts, extra)
    doc.save(path)


# ─────────────────────────────────────────────────────────────────────────────
# PowerPoint helpers
# ─────────────────────────────────────────────────────────────────────────────

def ppt_new() -> Presentation:
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    return prs


def ppt_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if slide.placeholders[1]:
        slide.placeholders[1].text = subtitle
    for ph in slide.placeholders:
        for para in ph.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = BRAND_BLUE_PPT


def ppt_section_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.text = bullets[0] if bullets else ""
    for b in bullets[1:]:
        p = body.add_paragraph()
        p.text = b
        p.level = 0
    for para in slide.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = BRAND_BLUE_PPT


def write_ppt(path: Path, title: str, subtitle: str, sections: list[tuple[str, list[str]]]) -> None:
    prs = ppt_new()
    ppt_title_slide(prs, title, subtitle)
    # Agenda
    ppt_section_slide(prs, "Agenda", [s[0] for s in sections])
    for s_title, bullets in sections:
        ppt_section_slide(prs, s_title, bullets)
    ppt_section_slide(prs, "Next steps", [
        "Confirm owners and dates",
        "Schedule follow-ups",
        "File this deck in the engagement folder",
    ])
    prs.save(path)


# ─────────────────────────────────────────────────────────────────────────────
# Excel helpers
# ─────────────────────────────────────────────────────────────────────────────

def xlsx_header(ws, headers: list[str]) -> None:
    for col, h in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col, value=h)
        cell.fill = HEADER_FILL
        cell.font = HEADER_FONT
        cell.alignment = Alignment(horizontal="left", vertical="center")
        cell.border = BORDER
        ws.column_dimensions[get_column_letter(col)].width = max(18, len(h) + 4)
    ws.freeze_panes = "A2"


def xlsx_status_dropdown(ws, col_letter: str, start_row: int = 2, end_row: int = 200) -> None:
    dv = DataValidation(type="list", formula1=STATUS_OPTIONS, allow_blank=True)
    dv.add(f"{col_letter}{start_row}:{col_letter}{end_row}")
    ws.add_data_validation(dv)


def xlsx_rows(ws, rows: list[list[str]], start_row: int = 2) -> None:
    for r_idx, row in enumerate(rows, start=start_row):
        for c_idx, val in enumerate(row, start=1):
            cell = ws.cell(row=r_idx, column=c_idx, value=val)
            cell.alignment = Alignment(vertical="top", wrap_text=True)
            cell.border = BORDER


# ─────────────────────────────────────────────────────────────────────────────
# Word artefacts
# ─────────────────────────────────────────────────────────────────────────────

def gen_qualification_questionnaire() -> Path:
    sections: list[tuple[str, list[str], list[str] | None]] = [
        ("Section 1 — Business context", [
            "What outcome does the AI app need to drive? (revenue, deflection, cycle time, quality)",
            "Who owns that outcome today, and what is the current baseline?",
            "What is the consequence of doing nothing for 12 months?",
        ], ["Score (0–3):", "Notes:"]),
        ("Section 2 — Authority & budget", [
            "Who is the executive sponsor? Are they on the kickoff call?",
            "Is there a budget line item, or is funding pending business case?",
            "Who signs the SOW?",
        ], ["Score (0–3):", "Notes:"]),
        ("Section 3 — Use case shape", [
            "RAG / grounded answers, agentic workflow, content generation, or extraction/classification?",
            "Single-tenant internal, multi-tenant SaaS, or customer-facing public endpoint?",
            "Synchronous (chat/API) or asynchronous (batch)?",
        ], ["Score (0–3):", "Notes:"]),
        ("Section 4 — Azure & data readiness", [
            "Existing Azure subscription? Landing zone in place?",
            "Are relevant data sources already in Azure?",
            "Azure OpenAI / AI Foundry provisioned in a supported region?",
            "Existing Entra ID + PIM posture?",
        ], ["Score (0–3):", "Notes:"]),
        ("Section 5 — Responsible AI posture", [
            "Is there an internal Responsible AI policy or standard?",
            "Has the use case been reviewed for high-risk categories?",
            "Who is the accountable human in the loop?",
        ], ["Score (0–3):", "Notes:"]),
        ("Section 6 — Operating model", [
            "Who operates the AI app after handover?",
            "Existing on-call rota?",
            "Model / prompt change-control expectation?",
        ], ["Score (0–3):", "Notes:"]),
        ("Overall scoring & gate", [
            "Total score (max 18):",
            "Gate: proceed only if total ≥ 12 AND no section scored 0.",
            "Go / no-go recommendation:",
            "Recommended next step (discovery workshop / remediation plan / disqualify):",
        ], None),
    ]
    p = OUT / "engagement" / "qualification-questionnaire.docx"
    write_word(p, "Qualification Questionnaire",
               "AI Applications on Microsoft Azure — engagement qualification", sections, "AI Apps presales lead")
    return p


def gen_definition_of_done() -> Path:
    sections = [
        ("Delivery artefacts", [
            "Signed HLD filed against the engagement",
            "Signed LLD filed against the engagement",
            "Runbook delivered and walked through with customer ops",
            "KT plan executed; sessions recorded and indexed",
            "Hypercare plan agreed, with start and end dates",
        ], None),
        ("Quality gates", [
            "WAF assessment re-run pre go-live; all Critical findings remediated or formally accepted",
            "Responsible AI checklist signed off by named accountable owner",
            "Eval dataset run on production model deployment; quality target met",
            "Load / performance test executed against agreed SLO",
            "Security review or pen test completed (where in scope)",
        ], None),
        ("Operational handover", [
            "Production access reduced to least privilege; engagement team off break-glass",
            "Alerts wired to customer's on-call; tested with synthetic incident",
            "Cost alerts configured against the agreed budget",
            "Prompt / model change-control process documented in the runbook",
            "Backup, restore, and DR validated (or formally scoped out)",
        ], None),
        ("Commercial closeout", [
            "Final invoice issued and acknowledged",
            "CSAT captured (target ≥ 8 / 10)",
            "Reference customer status agreed",
            "Case study draft (anonymised) within 30 days",
        ], None),
        ("Innersource closeout", [
            "At least one lesson-learned issue raised in the toolkit repo",
            "Any template improvement promoted to a PR",
            "Engagement folder archived per retention policy",
        ], None),
        ("Sign-off", [
            "Customer sponsor name + signature + date:",
            "Engagement lead name + signature + date:",
        ], None),
    ]
    p = OUT / "engagement" / "definition-of-done.docx"
    write_word(p, "Definition of Done",
               "AI Applications on Microsoft Azure — engagement closeout", sections, "Delivery quality lead")
    return p


def gen_hld() -> Path:
    sections = [
        ("Executive summary", [
            "Use case in one paragraph.",
            "Chosen reference architecture pattern (1/2/3/4) and why.",
            "Headline cost and timeline.",
        ], None),
        ("Business context & success metrics", [
            "Measurable change the AI app must drive (baseline → target).",
        ], None),
        ("Scope", [
            "In scope:",
            "Out of scope:",
            "Deferred to a later phase:",
        ], None),
        ("Solution architecture", [
            "Architecture block diagram (insert here).",
            "Narrative description.",
            "Chosen Azure services with justification vs. alternatives.",
        ], None),
        ("Data architecture", [
            "Sources, classification, flows, retention, residency.",
        ], None),
        ("Identity & access", [
            "Entra ID model, managed identities, RBAC scopes, conditional access.",
        ], None),
        ("Networking", [
            "Public vs. private, private endpoints, egress, DNS, WAF.",
        ], None),
        ("Responsible AI", [
            "Harms identification, mitigations, human-in-the-loop, transparency notes.",
        ], None),
        ("Non-functional requirements", [
            "Availability, latency, throughput, RPO/RTO, scalability, cost envelope.",
        ], None),
        ("Operations", [
            "Observability, alerting, on-call, runbook reference, change control.",
        ], None),
        ("Risks & assumptions", [
            "Register with owner and mitigation.",
        ], None),
        ("Phased delivery plan", [
            "Milestones, gates, deliverables per phase.",
        ], None),
        ("Sign-off", [
            "Customer sponsor:",
            "Engagement architect:",
        ], None),
    ]
    p = OUT / "deliverables" / "hld-template.docx"
    write_word(p, "High-Level Design (HLD)",
               "AI Applications on Microsoft Azure — engagement template", sections, "Engagement architect")
    return p


def gen_lld() -> Path:
    sections = [
        ("Resource inventory", [
            "Table of every Azure resource: name, SKU, region, redundancy, tags, owner.",
        ], None),
        ("Sizing & capacity", [
            "Model tokens-per-minute / RPM, vector index size, compute CPU/RAM, storage tiers, scale rules.",
        ], None),
        ("Network design", [
            "VNets, subnets, NSGs, private endpoints, route tables, DNS zones, ExpressRoute/VPN.",
        ], None),
        ("Identity design", [
            "Entra app registrations, managed identities, role assignments at exact scope, group memberships, CA policies.",
        ], None),
        ("Data design", [
            "Schemas, partitioning, indexing, retention, backup, encryption keys, Purview registrations.",
        ], None),
        ("AI model & prompt design", [
            "Azure OpenAI deployments (model, version, capacity, PTU vs. PAYG).",
            "Prompt repository structure, system prompts, eval harness, content-filter policy IDs.",
        ], None),
        ("Application configuration", [
            "Environment variables, Key Vault references, feature flags.",
        ], None),
        ("Observability", [
            "Application Insights resources, KQL queries for token usage, alert rules, dashboards, retention.",
        ], None),
        ("CI/CD", [
            "Repo layout, branching, GitHub Actions / Azure Pipelines stages, IaC tool, promotion gates, secret handling.",
        ], None),
        ("Security controls trace", [
            "Every WAF + RAI control mapped to the LLD section that implements it.",
        ], None),
        ("Test plan", [
            "Unit, integration, RAI eval, load, chaos, security.",
        ], None),
        ("Runbook reference", [
            "Pointer to the runbook deliverable.",
        ], None),
        ("Sign-off", [
            "Customer architect:",
            "Engagement architect:",
        ], None),
    ]
    p = OUT / "deliverables" / "lld-template.docx"
    write_word(p, "Low-Level Design (LLD)",
               "AI Applications on Microsoft Azure — engagement template", sections, "Engagement architect")
    return p


def gen_runbook() -> Path:
    sections = [
        ("At-a-glance card", [
            "App name, owners, prod region(s), SLO, key dashboards, key alerts, escalation path.",
        ], None),
        ("Architecture refresher", [
            "Diagram + 3-sentence narrative.",
        ], None),
        ("Routine operations", [
            "Daily/weekly health checks, backup verification, cost review, quota & capacity monitoring.",
        ], None),
        ("Incident response per alert", [
            "Per alert: trigger, confirm, mitigate, rollback, comms template.",
        ], None),
        ("Common AI-Apps incidents", [
            "Azure OpenAI 429 / quota exhaustion.",
            "Content filter false positives.",
            "Prompt injection signal in logs.",
            "RAG quality regression.",
            "Application Insights cost spike.",
        ], None),
        ("Change control", [
            "Prompt change: PR + two reviewers + eval set + canary 10% / 24h.",
            "Model change: full eval + content-filter retest + canary 10% / 7d + RAI sign-off.",
            "IaC change: standard CI/CD.",
        ], None),
        ("Disaster recovery", [
            "RPO/RTO, secondary region activation, data restore, failover test cadence.",
        ], None),
        ("Contacts", [
            "Customer ops, your hypercare team, Microsoft support (with support plan).",
        ], None),
    ]
    p = OUT / "deliverables" / "runbook-template.docx"
    write_word(p, "Operational Runbook",
               "AI Applications on Microsoft Azure — engagement template", sections, "SRE lead")
    return p


def gen_kt_plan() -> Path:
    sections = [
        ("Audience matrix", [
            "Per receiving role: current skill, required end-state, sessions assigned.",
        ], None),
        ("Session catalogue", [
            "Per session: title, duration, audience, prerequisites, learning objectives, hands-on exercise, owner, recorded yes/no.",
            "Standard sessions: architecture walkthrough; repo tour; prompt repo; eval harness; observability; cost & capacity; incident response; RAI controls; DR drill.",
        ], None),
        ("Hands-on labs", [
            "At least three sessions include a lab on customer infrastructure.",
        ], None),
        ("Exit assessment", [
            "Per individual: tasks they can perform unsupervised; signed.",
        ], None),
        ("Recording & artefact index", [
            "Every session recorded, indexed, linked from runbook.",
        ], None),
        ("Open-questions log", [
            "Questions raised + owner + resolution date.",
        ], None),
        ("Sign-off", [
            "Customer training lead:",
            "Engagement lead:",
        ], None),
    ]
    p = OUT / "deliverables" / "kt-plan-template.docx"
    write_word(p, "Knowledge Transfer Plan",
               "AI Applications on Microsoft Azure — engagement template", sections, "Delivery lead")
    return p


def gen_hypercare_plan() -> Path:
    sections = [
        ("Duration & window", [
            "Start date, end date, daily working hours, time zone, weekend coverage.",
        ], None),
        ("Team & rota", [
            "Your engineers, customer engineers shadowing, escalation tree, Microsoft support plan.",
        ], None),
        ("Daily cadence", [
            "15-min standup, KPI snapshot, action log update.",
        ], None),
        ("Weekly cadence", [
            "1-hour steerco: KPIs vs. SLO, open incidents, change requests, exit-gate readiness.",
        ], None),
        ("KPIs to track", [
            "Availability vs. SLO, p95 latency, RAG/model quality on eval set, content filter / abuse signals, cost vs. budget, incident count by severity.",
        ], None),
        ("Incident handling", [
            "Tighter response-time SLAs than BAU; align to runbook.",
        ], None),
        ("Change control", [
            "Restricted: only blocking fixes; document exceptions process.",
        ], None),
        ("Exit gate", [
            "14 consecutive days without a Sev-1, all Sev-2s closed, SLOs met, customer on-call handled at least one real incident unaided.",
        ], None),
        ("Handover at exit", [
            "Formal sign-off, BAU transition, retrospective with lessons learned filed.",
        ], None),
        ("Sign-off", [
            "Customer sponsor:",
            "Engagement lead:",
        ], None),
    ]
    p = OUT / "deliverables" / "hypercare-plan-template.docx"
    write_word(p, "Hypercare Plan",
               "AI Applications on Microsoft Azure — engagement template", sections, "Delivery lead")
    return p


# ─────────────────────────────────────────────────────────────────────────────
# PowerPoint artefacts
# ─────────────────────────────────────────────────────────────────────────────

def gen_offering_one_pager() -> Path:
    sections = [
        ("Headline value proposition", [
            "<One sentence: who you help, to do what, with which Azure AI services>",
        ]),
        ("Packages", [
            "Envisioning Workshop — 1 week",
            "RAG Pilot on Azure OpenAI / AI Foundry — 4 weeks",
            "Production AI App on AKS / Azure Container Apps — 8–12 weeks",
        ]),
        ("Outcomes", [
            "Live RAG chatbot grounded on customer documents",
            "Measured answer quality > 80% on a 50-question eval set",
            "Production-grade deployment with Azure Monitor + RAI guardrails",
        ]),
        ("Azure stack", [
            "Azure OpenAI in Azure AI Foundry",
            "Azure AI Search (vector + hybrid)",
            "Azure Container Apps / AKS / App Service",
            "Azure Cosmos DB / Azure SQL / Microsoft Fabric",
            "Azure Monitor + Application Insights",
            "Microsoft Entra ID",
        ]),
        ("Responsible AI commitment", [
            "Mapped to Microsoft Responsible AI Standard",
            "Content filters, abuse monitoring, human-in-the-loop where applicable",
        ]),
        ("Proof points", [
            "Customer A — <industry, outcome>",
            "Customer B — <industry, outcome>",
            "Customer C — <industry, outcome>",
        ]),
        ("Call to action", [
            "Book a 60-minute envisioning workshop",
            "Contact: <named contact>",
        ]),
    ]
    p = OUT / "engagement" / "offering-one-pager.pptx"
    write_ppt(p, "AI Applications on Azure", "Offering one-pager (template)", sections)
    return p


def gen_discovery_deck() -> Path:
    sections = [
        ("Welcome, scope & success criteria", [
            "Workshop objectives",
            "Roles in the room",
            "Definition of success for the two days",
        ]),
        ("Day 1 — Use case framing", [
            "Walk the candidate use cases",
            "Score: value × feasibility × risk",
            "Pick top 1–2",
        ]),
        ("Day 1 — Data strategy", [
            "Inventory data sources",
            "Capture location, sensitivity, refresh, owner, lineage",
            "Flag non-Azure-resident sources",
        ]),
        ("Day 2 — Architecture & service selection", [
            "Map use case to a reference architecture pattern",
            "Choose compute (ACA / AKS / App Service)",
            "Choose data (Cosmos DB / Azure SQL / Fabric)",
        ]),
        ("Day 2 — Responsible AI & operations", [
            "Walk the Microsoft Responsible AI Standard",
            "Human-in-the-loop checkpoints",
            "Content filters, abuse monitoring",
            "Change control for prompts and models",
        ]),
        ("Closeout", [
            "Confirm next deliverable (HLD + costed proposal)",
            "Agree delivery date",
            "Capture lessons & open questions",
        ]),
    ]
    p = OUT / "engagement" / "discovery-workshop-deck.pptx"
    write_ppt(p, "AI Apps Discovery Workshop", "2-day facilitator deck (template)", sections)
    return p


# ─────────────────────────────────────────────────────────────────────────────
# Excel artefacts
# ─────────────────────────────────────────────────────────────────────────────

def gen_discovery_workbook() -> Path:
    wb = Workbook()
    # Tab 1 — Use cases
    ws = wb.active
    ws.title = "1. Use cases"
    xlsx_header(ws, ["#", "Use case", "Sponsor", "Value (1–5)", "Feasibility (1–5)", "Risk (1–5)", "Score", "Selected?", "Notes"])
    xlsx_rows(ws, [
        ["1", "<RAG over policy documents>", "<Name>", "", "", "", "", "No", ""],
        ["2", "<Agent for service desk>", "<Name>", "", "", "", "", "No", ""],
    ])
    xlsx_status_dropdown(ws, "H")

    # Tab 2 — Data inventory
    ws2 = wb.create_sheet("2. Data inventory")
    xlsx_header(ws2, ["#", "Source name", "Location", "Sensitivity", "Refresh cadence", "Owner", "In Azure?", "Notes"])
    xlsx_rows(ws2, [
        ["1", "<SharePoint policies>", "<URL>", "Confidential", "Weekly", "<Owner>", "Yes", ""],
        ["2", "<CRM extract>", "<System>", "Restricted", "Daily", "<Owner>", "No", "Needs Synapse pipeline"],
    ])

    # Tab 3 — Architecture decisions
    ws3 = wb.create_sheet("3. Architecture decisions")
    xlsx_header(ws3, ["Decision area", "Options considered", "Selected", "Rationale", "Owner"])
    xlsx_rows(ws3, [
        ["Reference pattern", "RAG / Agent / Real-time / Fabric", "RAG", "Fits document-grounding use case", "<Architect>"],
        ["Compute tier", "ACA / AKS / App Service", "Azure Container Apps", "No existing AKS investment", "<Architect>"],
        ["Data tier", "Cosmos DB / Azure SQL / Fabric", "Cosmos DB for NoSQL", "Conversation history with TTL", "<Architect>"],
        ["Embedding model", "text-embedding-3-small / -large", "text-embedding-3-large", "Accuracy on long-form policy text", "<Architect>"],
    ])

    # Tab 4 — Responsible AI checklist
    ws4 = wb.create_sheet("4. Responsible AI")
    xlsx_header(ws4, ["#", "Check", "Status", "Owner", "Evidence"])
    rai_checks = [
        "Use case reviewed for high-risk categories",
        "Harms identified and documented",
        "Content safety policy selected",
        "Abuse monitoring enabled",
        "Human-in-the-loop checkpoints defined",
        "Transparency note authored",
        "Eval set scoped",
        "Post-deployment monitoring plan agreed",
    ]
    xlsx_rows(ws4, [[str(i+1), c, "Not started", "<Owner>", ""] for i, c in enumerate(rai_checks)])
    xlsx_status_dropdown(ws4, "C")

    # Tab 5 — Next steps log
    ws5 = wb.create_sheet("5. Next steps")
    xlsx_header(ws5, ["#", "Action", "Owner", "Due", "Status", "Notes"])
    xlsx_rows(ws5, [["1", "Draft HLD", "<Architect>", "yyyy-mm-dd", "Not started", ""]])
    xlsx_status_dropdown(ws5, "E")

    p = OUT / "engagement" / "discovery-workshop-workbook.xlsx"
    wb.save(p)
    return p


def gen_waf_workbook() -> Path:
    wb = Workbook()
    pillars = [
        ("Reliability", [
            ("Model availability & quota headroom", "Azure OpenAI quota dashboard"),
            ("Secondary deployment in paired region", "Resource Graph inventory"),
            ("SDK retry + fallback strategy", "Code review"),
            ("Graceful degradation when LLM unavailable", "Failure-mode test"),
        ]),
        ("Security", [
            ("Entra ID auth on Azure OpenAI (no key auth)", "RBAC review"),
            ("Private endpoints for AI Services + AI Search", "Network diagram"),
            ("Managed identity end-to-end", "Identity review"),
            ("Azure AI Content Safety configured", "Policy ID + test prompts"),
            ("Prompt-injection mitigation", "System prompt + filters"),
            ("Data-exfiltration controls (egress)", "NSG / Firewall rules"),
        ]),
        ("Cost Optimization", [
            ("Right model for task (mini vs. flagship)", "Eval results across models"),
            ("PTU vs. PAYG decision", "Cost model"),
            ("Prompt-size discipline", "Token telemetry"),
            ("Semantic caching", "Cache hit rate"),
            ("Per-tenant token budgets", "Quota policy"),
        ]),
        ("Operational Excellence", [
            ("Prompt & response telemetry (redacted)", "Application Insights"),
            ("Eval pipeline in CI", "GitHub Actions log"),
            ("Prompt change control (PR + canary)", "Repo settings"),
            ("Model rollout strategy (canary/shadow)", "Deployment runbook"),
            ("Cost & quota alerts", "Azure Monitor"),
        ]),
        ("Performance Efficiency", [
            ("Latency budget per request defined", "SLO doc"),
            ("Streaming responses end-to-end", "Code review"),
            ("Vector index tuning (HNSW params)", "AI Search settings"),
            ("AI Search replicas / partitions sized", "AI Search settings"),
            ("Parallelism in orchestration", "Code review"),
        ]),
        ("Responsible AI", [
            ("Harms identification documented", "RAI assessment doc"),
            ("Content filter policy approved", "RAI sign-off"),
            ("Abuse monitoring enabled", "Azure portal"),
            ("Human-in-the-loop where required", "Workflow design"),
            ("Transparency note published", "Doc link"),
            ("Post-deployment monitoring against eval set", "Eval dashboard"),
        ]),
    ]
    first = True
    for pillar, checks in pillars:
        ws = wb.active if first else wb.create_sheet(pillar)
        if first:
            ws.title = pillar
            first = False
        xlsx_header(ws, ["#", "Check", "Severity if missing", "Status", "Owner", "Finding ID", "Evidence / link", "Remediation notes"])
        rows = [[str(i+1), c, "TBD", "Not started", "<Owner>", "", e, ""] for i, (c, e) in enumerate(checks)]
        xlsx_rows(ws, rows)
        xlsx_status_dropdown(ws, "D")
        sev_dv = DataValidation(type="list", formula1='"Critical,High,Medium,Low"', allow_blank=True)
        sev_dv.add(f"C2:C200")
        ws.add_data_validation(sev_dv)

    # Findings tab
    ws_f = wb.create_sheet("Findings")
    xlsx_header(ws_f, ["ID", "Pillar", "Finding", "Severity", "Owner", "Status", "Target fix date", "Remediation"])
    xlsx_rows(ws_f, [
        ["WAF-001", "Security", "<Finding>", "High", "<Owner>", "Open", "yyyy-mm-dd", "<Remediation>"],
    ])
    xlsx_status_dropdown(ws_f, "F")

    p = OUT / "engagement" / "waf-assessment.xlsx"
    wb.save(p)
    return p


def gen_assessment_inputs() -> Path:
    wb = Workbook()
    ws = wb.active
    ws.title = "Inputs checklist"
    xlsx_header(ws, ["#", "Input", "Source", "Format", "Provided?", "Owner", "Notes"])
    inputs = [
        ("Azure tenant ID and primary subscription IDs", "Azure portal", "Text"),
        ("Landing zone state (none / partial / ALZ-aligned)", "Customer architect", "Document"),
        ("Azure Monitor / Log Analytics workspace IDs", "Azure portal", "Text"),
        ("Application Insights resource IDs for in-scope apps", "Azure portal", "Text"),
        ("Azure OpenAI / AI Foundry resource list + region(s)", "Azure portal", "Spreadsheet"),
        ("Model deployments (model, version, capacity)", "Azure OpenAI / Foundry", "Spreadsheet"),
        ("Existing Azure spend by service (last 90 days)", "Cost Management export", "CSV"),
        ("Quota requests history and current consumption", "Quota dashboard", "Export"),
        ("Identity model (Entra tenant, B2B/B2C, PIM state)", "Customer identity team", "Document"),
        ("Data sources in scope (location, classification, owner)", "Customer data team", "Spreadsheet"),
        ("Connectivity (ExpressRoute / VPN / public)", "Customer network team", "Diagram"),
        ("Current Responsible AI governance posture", "Customer compliance", "Document"),
    ]
    rows = [[str(i+1), name, src, fmt, "No", "<Owner>", ""] for i, (name, src, fmt) in enumerate(inputs)]
    xlsx_rows(ws, rows)
    dv = DataValidation(type="list", formula1='"Yes,No,Partial"', allow_blank=True)
    dv.add("E2:E200")
    ws.add_data_validation(dv)

    ws2 = wb.create_sheet("Azure inventory (ARG)")
    xlsx_header(ws2, ["Resource type", "Name", "Location", "Resource group", "Subscription", "Tags", "Notes"])
    xlsx_rows(ws2, [["Microsoft.CognitiveServices/accounts", "<aoai-prod>", "swedencentral", "<rg>", "<sub>", "env=prod", ""]])

    p = OUT / "engagement" / "assessment-platform-inputs.xlsx"
    wb.save(p)
    return p


def gen_evidence_tracker() -> Path:
    wb = Workbook()
    ws = wb.active
    ws.title = "Audit evidence tracker"
    xlsx_header(ws, ["Module", "Control", "Topic", "Evidence item", "Owner", "Status", "Last updated", "File name", "Notes"])
    rows: list[list[str]] = []
    module_a = [
        ("A.1.1", "Organisational Data", ["Certificate of incorporation", "Org chart", "Key personnel list"]),
        ("A.1.2", "Financial Documentation", ["Financial statements", "Professional indemnity insurance"]),
        ("A.2.1", "Service Delivery Methodology", ["Delivery playbook", "SOW template", "Kickoff artefact", "Closure document", "Risk register template"]),
        ("A.2.2", "Quality Management", ["QMS policy", "CSAT process", "Escalation procedure"]),
        ("A.3.1", "Customer Satisfaction Outcomes", ["CSAT/NPS data", "References", "Testimonials"]),
        ("A.3.2", "Complaint Handling", ["Complaint register", "Resolved case", "Root cause analysis"]),
        ("A.3.3", "Security & Privacy", ["InfoSec policy", "Data protection / GDPR", "Breach procedure", "Staff training records"]),
    ]
    module_b = [
        ("B.1.1", "Azure AI Implementation Capability", ["Case studies", "Architecture diagrams", "Capability statement"]),
        ("B.2.1", "ACR Performance", ["ACR export Pillar 1", "ACR export Pillar 2", "ACR export Pillar 3"]),
        ("B.2.2", "Customer Diversity", ["DPOR/PAL/CSP report (≥3 customers)"]),
        ("B.3.1", "Certifications", ["AZ-204 transcripts", "AZ-400 transcripts", "AI-102 transcripts", "DP-420 transcripts", "Roster of ≥5 individuals"]),
        ("B.4.1", "Audit Readiness", ["Evidence index", "Internal review log", "Submission package"]),
        ("B.4.2", "Partner Onboarding Assets", ["Customer onboarding pack", "Delivery templates", "KT plan", "Runbook"]),
    ]
    for module, items in (("Module A", module_a), ("Module B", module_b)):
        for ref, topic, evs in items:
            for ev in evs:
                rows.append([module, ref, topic, ev, "<Owner>", "Not started", "yyyy-mm-dd", "", ""])
    xlsx_rows(ws, rows)
    xlsx_status_dropdown(ws, "F", end_row=len(rows) + 1)

    # Summary tab
    ws2 = wb.create_sheet("Summary by control")
    xlsx_header(ws2, ["Module", "Control", "Topic", "Items", "Complete", "% complete"])
    summary: list[list[str]] = []
    for module, items in (("Module A", module_a), ("Module B", module_b)):
        for ref, topic, evs in items:
            summary.append([module, ref, topic, str(len(evs)), "0", "0%"])
    xlsx_rows(ws2, summary)

    p = OUT / "audit" / "evidence-tracker.xlsx"
    wb.save(p)
    return p


def gen_prequal_checklist() -> Path:
    wb = Workbook()
    ws = wb.active
    ws.title = "Pre-qualification gate"
    xlsx_header(ws, ["#", "Requirement", "Detail", "Status", "Evidence", "Owner", "Notes"])
    rows = [
        ["1", "Solutions Partner designation", "Data & AI (Azure) OR Digital & App Innovation (Azure) — active in Partner Center", "Not started", "", "<Owner>", ""],
        ["2", "ACR Pillar 1 – AI Services", "≥ $15,000 USD last 3 months: Azure OpenAI, Rest of Azure AI, 3P GPU, Microsoft Foundry", "Not started", "", "<Owner>", ""],
        ["3", "ACR Pillar 2 – App Platform", "≥ $15,000 USD last 3 months: AKS, ACA, ARO, App Service, Logic Apps, APIM, Functions, Managed Redis, GitHub", "Not started", "", "<Owner>", ""],
        ["4", "ACR Pillar 3 – Data Platform", "≥ $15,000 USD last 3 months: Cosmos DB, SQL Hyperscale, Azure SQL Core, MySQL PaaS, PostgreSQL PaaS, Fabric F SKU", "Not started", "", "<Owner>", ""],
        ["5", "Customer diversity", "≥ 3 unique customers via DPOR, PAL, or CSP", "Not started", "", "<Owner>", ""],
        ["6", "Certifications", "≥ 5 individuals; each of AZ-204, AZ-400, AI-102, DP-420 held by ≥ 1 person", "Not started", "", "<Owner>", ""],
    ]
    xlsx_rows(ws, rows)
    xlsx_status_dropdown(ws, "D", end_row=len(rows) + 1)
    p = OUT / "audit" / "pre-qual-checklist.xlsx"
    wb.save(p)
    return p


# ─────────────────────────────────────────────────────────────────────────────
# Driver
# ─────────────────────────────────────────────────────────────────────────────

def main() -> None:
    ensure_dirs()
    paths = [
        gen_offering_one_pager(),
        gen_qualification_questionnaire(),
        gen_discovery_deck(),
        gen_discovery_workbook(),
        gen_waf_workbook(),
        gen_assessment_inputs(),
        gen_definition_of_done(),
        gen_hld(),
        gen_lld(),
        gen_runbook(),
        gen_kt_plan(),
        gen_hypercare_plan(),
        gen_evidence_tracker(),
        gen_prequal_checklist(),
    ]
    for p in paths:
        rel = p.relative_to(REPO_ROOT)
        print(f"wrote {rel}")
    print(f"\n{len(paths)} workfiles generated under {OUT.relative_to(REPO_ROOT)}/")


if __name__ == "__main__":
    main()
